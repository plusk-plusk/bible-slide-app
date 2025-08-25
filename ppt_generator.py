from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.oxml.xmlchemy import OxmlElement   # ✅ 텍스트 윤곽선(Outline) 설정용
from extractor import extract_bible_verses

# ==============================
# 튜닝용 상수 (필요 시 숫자만 조정)
# ==============================
# 초미세 그림자(패널 없음)
SHADOW_COLOR     = RGBColor(22, 24, 20)  # 검정보다 부드러운 짙은 올리브
SHADOW_OFFSET_IN = 0.02                  # 0.015 ~ 0.03 권장

# 텍스트 윤곽선(Outline)
OUTLINE_HEX   = "222222"  # 짙은 회색(16진 RGB). 더 부드럽게는 "2A2A2A" 등
OUTLINE_PT    = 1.0       # 윤곽선 두께(pt). 0.75~1.25 권장

# 폰트 크기/줄간격(요청대로 유지)
TITLE_SIZE_PT = 36
BODY_SIZE_PT  = 50
LINE_SPACING  = 1.15

# 레이아웃(여백만 소폭 활용)
OUTER_MARGIN_IN   = 1.3
INNER_SIDE_GAP_IN = 0.25

# 선택: 배경 전체를 아주 살짝만 소프닝(밝은 흰막 90% 투명)
USE_GLOBAL_SOFTEN_OVERLAY = False   # 필요하면 True
SOFTEN_ALPHA = 0.90                 # 0.88~0.92 정도가 안전

# ==============================
# 유틸 함수
# ==============================
def apply_text_outline(run, color_hex="222222", width_pt=1.0):
    """
    텍스트 Run에 윤곽선(Outline) 추가.
    - color_hex: "RRGGBB" (예: "222222")
    - width_pt : pt 단위 두께 (1pt ≈ 12700 DrawingML 단위)
    """
    r = run._r
    rPr = r.get_or_add_rPr()

    # <a:ln w="..."><a:solidFill><a:srgbClr val="..."/></a:solidFill></a:ln>
    ln = OxmlElement('a:ln')
    ln.set('w', str(int(width_pt * 12700)))  # 약 1pt

    solidFill = OxmlElement('a:solidFill')
    srgbClr = OxmlElement('a:srgbClr')
    srgbClr.set('val', color_hex)
    solidFill.append(srgbClr)
    ln.append(solidFill)

    # 기존에 ln이 있다면 교체, 없으면 추가
    # (간단화: 그냥 append, 다중 추가되면 마지막 것이 적용됨)
    rPr.append(ln)

def add_text_with_micro_shadow(slide, text, left, top, width, height,
                               font_size_pt=48, bold=True, align=PP_ALIGN.LEFT,
                               font_name='Apple SD Gothic Neo',
                               with_outline=True):
    """
    패널 없이 텍스트만 두 겹(뒤: 그림자 / 앞: 흰색)으로 배치.
    - 그림자 오프셋을 매우 작게 해서 '은은한 윤곽'만 남김
    - 흰 텍스트에 '윤곽선(Outline)' 추가로 대비 강화 (노년층 가독성↑)
    """
    # 1) 그림자 층
    s_off = Inches(SHADOW_OFFSET_IN)
    shadow_box = slide.shapes.add_textbox(left + s_off, top + s_off, width, height)
    sframe = shadow_box.text_frame
    sframe.word_wrap = True
    p = sframe.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size_pt)
    run.font.bold = bold
    run.font.color.rgb = SHADOW_COLOR
    run.font.name = font_name

    # 2) 실제 텍스트(앞쪽)
    main_box = slide.shapes.add_textbox(left, top, width, height)
    mframe = main_box.text_frame
    mframe.word_wrap = True
    p2 = mframe.paragraphs[0]
    p2.alignment = align
    run2 = p2.add_run()
    run2.text = text
    run2.font.size = Pt(font_size_pt)
    run2.font.bold = bold
    run2.font.color.rgb = RGBColor(255, 255, 255)  # 흰색
    run2.font.name = font_name

    # ✅ 윤곽선(Outline)로 대비 강화
    if with_outline:
        apply_text_outline(run2, color_hex=OUTLINE_HEX, width_pt=OUTLINE_PT)

    return main_box, shadow_box

def add_soften_overlay(slide, prs, alpha=0.90):
    """배경을 아주 살짝만 '흰막'으로 소프닝(글자 대비↑, 색감은 최대한 유지)"""
    from pptx.enum.shapes import MSO_SHAPE
    ov = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    fill = ov.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(255, 255, 255)
    fill.transparency = alpha
    ov.line.fill.background()
    return ov

# ==============================
# 메인
# ==============================
def make_bible_ppt(json_path, ref_path, output_path, background_image):
    prs = Presentation()
    prs.slide_width  = Inches(13.33)  # 16:9
    prs.slide_height = Inches(7.5)

    verses = extract_bible_verses(json_path, ref_path)
    if not verses:
        print("⚠️ 구절이 없어서 PPT를 생성하지 않았습니다.")
        return

    for verse in verses:
        slide = prs.slides.add_slide(prs.slide_layouts[6])

        # 배경
        slide.shapes.add_picture(
            background_image, 0, 0,
            width=prs.slide_width, height=prs.slide_height
        )

        # 선택: 배경 전체를 아주 살짝만 소프닝(노년층 가독성↑)
        if USE_GLOBAL_SOFTEN_OVERLAY:
            add_soften_overlay(slide, prs, alpha=SOFTEN_ALPHA)

        # 레이아웃(여백/폭)
        margin = Inches(OUTER_MARGIN_IN)
        left   = margin + Inches(INNER_SIDE_GAP_IN)
        right  = margin + Inches(INNER_SIDE_GAP_IN)
        box_w  = prs.slide_width - (left + right)

        # ----- 제목 -----
        title_top = Inches(1.0)
        title_box, title_shadow = add_text_with_micro_shadow(
            slide,
            text=verse["title"],
            left=left, top=title_top,
            width=box_w, height=Inches(0.9),
            font_size_pt=TITLE_SIZE_PT,
            bold=True, align=PP_ALIGN.LEFT,
            with_outline=True
        )

        # ----- 본문 -----
        body_top = title_top + Inches(0.70)
        body_h   = Inches(2.6)  # 높이는 기존과 동일 유지(원하면 배경에 맞춰 2.4~3.0 조정)

        main_box, shadow_box = add_text_with_micro_shadow(
            slide,
            text=verse["text"],
            left=left, top=body_top,
            width=box_w, height=body_h,
            font_size_pt=BODY_SIZE_PT,
            bold=True, align=PP_ALIGN.JUSTIFY,
            with_outline=True
        )

        # 줄간격 (요청대로 유지)
        main_box.text_frame.paragraphs[0].line_spacing = LINE_SPACING
        shadow_box.text_frame.paragraphs[0].line_spacing = LINE_SPACING

    prs.save(output_path)
    print(f"✅ PPT 저장 완료: {output_path}")
