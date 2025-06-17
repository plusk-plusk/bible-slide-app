from flask import Flask, request, send_file, render_template, redirect, url_for
from pptx import Presentation
from pptx.util import Inches
import json

app = Flask(__name__)

# 루트 경로 → /upload로 리디렉션
@app.route("/")
def index():
    return redirect(url_for("upload"))

# 구절 입력 폼 페이지
@app.route("/upload", methods=["GET"])
def upload():
    return render_template("index.html")

# 슬라이드 생성 요청
@app.route("/generate", methods=["POST"])
def generate():
    with open("bible.json", "r", encoding="utf-8") as f:
        bible_data = json.load(f)

    verses_raw = request.form.get("verses")
    if not verses_raw:
        return "❌ 구절을 입력해주세요.", 400

    verses = verses_raw.splitlines()
    prs = Presentation()
    blank_slide_layout = prs.slide_layouts[6]

    for verse in verses:
        verse = verse.strip()
        if not verse:
            continue
        try:
            book, chap_vers = verse.split()
            chapter, verse_number = map(int, chap_vers.split(":"))
            verse_text = bible_data[book][str(chapter)][str(verse_number)]
        except Exception as e:
            verse_text = f"⚠️ 구절 파싱 오류: {verse}"

        slide = prs.slides.add_slide(blank_slide_layout)
        textbox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(5.5))
        text_frame = textbox.text_frame
        text_frame.text = verse_text

    output_path = "BibleSlides.pptx"
    prs.save(output_path)
    return send_file(output_path, as_attachment=True)

# 선택적으로 /download로 슬라이드 바로 다운로드 가능
@app.route("/download")
def download():
    return send_file("BibleSlides.pptx", as_attachment=True)

if __name__ == "__main__":
    app.run(debug=True)
